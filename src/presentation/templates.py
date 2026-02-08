"""
Genesis Presentation Engine — Slide Templates (v2: Next-Gen)

AI-generated backgrounds, glassmorphism cards, glow accents,
dark-mode-aware rendering. The future of financial presentations.
"""

import os
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
from design import rgb, is_dark, is_editorial, SECTION_BADGE_COLORS, EDITORIAL_LAYOUT


# ============================================================================
# Low-Level Visual Primitives
# ============================================================================

def _set_shape_alpha(shape, alpha_pct: int):
    """
    Set fill transparency on a shape (0=transparent, 100=opaque).
    Uses direct XML manipulation on the shape's spPr element.
    """
    spPr = shape._element.spPr
    solid = spPr.find(qn('a:solidFill'))
    if solid is None:
        return
    clr = solid.find(qn('a:srgbClr'))
    if clr is None:
        return
    # Remove existing alpha if any
    for existing in clr.findall(qn('a:alpha')):
        clr.remove(existing)
    alpha = etree.SubElement(clr, qn('a:alpha'))
    alpha.set('val', str(alpha_pct * 1000))  # PowerPoint uses 0-100000


def add_slide_bg(slide, prs, palette, bg_image=None):
    """Add background to slide — AI image or solid color."""
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    if bg_image and os.path.exists(bg_image):
        slide.shapes.add_picture(bg_image, Emu(0), Emu(0), slide_w, slide_h)
    else:
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), slide_w, slide_h
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = rgb(palette.slide_bg)
        bg.line.fill.background()


def add_glass_card(slide, left, top, width, height, palette, alpha_pct=75):
    """
    Semi-transparent glassmorphism card.
    Creates a frosted-glass rectangle for overlaying on AI backgrounds.
    """
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = rgb(palette.card_bg)
    card.line.color.rgb = rgb(palette.card_border)
    card.line.width = Pt(0.5)
    _set_shape_alpha(card, alpha_pct)
    return card


def add_glow_line(slide, left, top, width, palette, color=None):
    """Accent glow line — a bright thin line for visual separation."""
    c = color or palette.gold
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Pt(2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(c)
    line.line.fill.background()
    return line


def add_hyperlink_run(paragraph, text, url, font_size=10, color="#00D4FF", bold=False, italic=False):
    """
    Add a clickable hyperlink run to a paragraph.
    Opens URL when clicked in presentation mode.
    """
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = rgb(color)
    run.font.bold = bold
    run.font.italic = italic
    run.font.underline = True
    if url:
        run.hyperlink.address = url
    return run


def add_video_button(slide, url, label="▶ Watch Video", left=10.5, top=0.7,
                     width=2.2, height=0.4, palette=None):
    """
    Add a clickable video link button — opens video URL in browser.
    Styled as a pill-shaped button with play icon.
    """
    btn = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    btn.fill.solid()
    btn.fill.fore_color.rgb = rgb(palette.red if palette else "#CC0000")
    btn.line.fill.background()

    # Add hyperlink to the shape
    btn.click_action.hyperlink.address = url

    tf = btn.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(2)
    r = p.add_run()
    r.text = label
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = rgb("#FFFFFF")
    r.font.name = "Arial"
    return btn


def add_link_button(slide, url, label="🔗 Open Link", left=10.5, top=0.7,
                    width=2.2, height=0.35, color=None, palette=None):
    """
    Generic clickable link button — opens URL in browser during presentation.
    """
    c = color or (palette.chart_primary if palette else "#003366")
    btn = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    btn.fill.solid()
    btn.fill.fore_color.rgb = rgb(c)
    btn.line.fill.background()
    btn.click_action.hyperlink.address = url

    tf = btn.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(1)
    r = p.add_run()
    r.text = label
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = rgb("#FFFFFF")
    r.font.name = "Arial"
    return btn


# ============================================================================
# Editorial Primitives (SYZ-style)
# ============================================================================

def add_editorial_header(slide, palette, meta: dict):
    """SYZ-style thin header: orange accent line + hashtag tag left + date right."""
    el = EDITORIAL_LAYOUT
    slide_w_in = meta.get("slide_width", 13.333)

    # 2pt orange accent line at very top
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(el.accent_line_top),
        Inches(slide_w_in), Inches(el.accent_line_height)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb(palette.orange)
    accent.line.fill.background()

    # Left: header tag in orange
    tag_box = slide.shapes.add_textbox(
        Inches(el.margin_left), Inches(el.header_tag_top),
        Inches(7), Inches(el.header_tag_height)
    )
    tf = tag_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.space_before = Pt(4)
    run = p.add_run()
    run.text = meta.get("header_tag", "#GLOBALMARKETS WEEKLY WRAP-UP")
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = rgb(palette.orange)
    run.font.name = "Arial"

    # Right: date
    date_box = slide.shapes.add_textbox(
        Inches(slide_w_in - 3.6), Inches(el.header_tag_top),
        Inches(3.0), Inches(el.header_tag_height)
    )
    tf2 = date_box.text_frame
    tf2.word_wrap = False
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    p2.space_before = Pt(4)
    run2 = p2.add_run()
    run2.text = meta.get("date", "")
    run2.font.size = Pt(10)
    run2.font.color.rgb = rgb(palette.gray)
    run2.font.name = "Arial"

    # Thin gray separator line
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(el.margin_left), Inches(el.header_sep_top),
        Inches(slide_w_in - el.margin_left - el.margin_right), Pt(1)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = rgb(palette.light_gray)
    sep.line.fill.background()


def add_section_badge(slide, section: str, left: float, top: float, palette):
    """Colored rounded-rect badge with white text (e.g., green '#equities')."""
    color = SECTION_BADGE_COLORS.get(section, palette.gray)
    label = f"#{section.replace('_', ' ')}"
    # Auto-width based on text length
    badge_w = max(1.2, len(label) * 0.11 + 0.3)

    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(badge_w), Inches(0.28)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = rgb(color)
    badge.line.fill.background()

    tf = badge.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = rgb("#FFFFFF")
    run.font.name = "Arial"
    return badge


def add_editorial_footer(slide, palette, meta: dict, page_num: int):
    """Editorial footer: logo/brand left + page number right, thin separator above."""
    el = EDITORIAL_LAYOUT
    slide_w_in = meta.get("slide_width", 13.333)

    # Thin separator line
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(el.margin_left), Inches(el.footer_sep_top),
        Inches(slide_w_in - el.margin_left - el.margin_right), Pt(1)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = rgb(palette.light_gray)
    sep.line.fill.background()

    # Left: logo image or text fallback
    logo_path = meta.get("logo_path", "")
    if logo_path and os.path.exists(logo_path):
        slide.shapes.add_picture(
            logo_path, Inches(el.margin_left), Inches(el.footer_top),
            Inches(1.8), Inches(0.3)
        )
    else:
        lb = slide.shapes.add_textbox(
            Inches(el.margin_left), Inches(el.footer_top),
            Inches(4), Inches(el.footer_height)
        )
        tf = lb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = meta.get("footer_left", meta.get("company", "Rossignoli & Partners"))
        run.font.size = Pt(8)
        run.font.bold = True
        run.font.color.rgb = rgb(palette.orange)
        run.font.name = "Arial"

    # Right: page number
    rb = slide.shapes.add_textbox(
        Inches(slide_w_in - 3.0), Inches(el.footer_top),
        Inches(2.4), Inches(el.footer_height)
    )
    tf2 = rb.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run2 = p2.add_run()
    run2.text = f"Page {page_num}"
    run2.font.size = Pt(8)
    run2.font.color.rgb = rgb(palette.source_color)
    run2.font.name = "Arial"


# ============================================================================
# Reusable Slide Furniture
# ============================================================================

def add_header_bar(slide, palette, meta: dict):
    """Navy bar at top with branding — works on both light and dark."""
    slide_w = Inches(meta.get("slide_width", 13.333))

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), slide_w, Inches(0.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(palette.navy)
    bar.line.fill.background()

    # Left: tag text
    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "   " + meta.get("header_tag", "#WEEKLY STRATEGY")
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = rgb(palette.gold)

    # Center: company name
    center_box = slide.shapes.add_textbox(
        Inches(4.5), Inches(0), Inches(4.333), Inches(0.5)
    )
    tf2 = center_box.text_frame
    tf2.word_wrap = False
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(6)
    run2 = p2.add_run()
    run2.text = meta.get("company", "Crossinvest SA")
    run2.font.size = Pt(11)
    run2.font.bold = True
    run2.font.color.rgb = rgb(palette.white)

    # Right: date
    right_box = slide.shapes.add_textbox(
        Inches(10), Inches(0), Inches(3.0), Inches(0.5)
    )
    tf3 = right_box.text_frame
    tf3.word_wrap = False
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.RIGHT
    p3.space_before = Pt(6)
    run3 = p3.add_run()
    run3.text = meta.get("date", "") + "   "
    run3.font.size = Pt(10)
    run3.font.color.rgb = rgb(palette.white)


def add_section_line(slide, palette):
    """Short accent line below header — gold on dark, navy on light."""
    dark = is_dark(palette)
    color = palette.gold if dark else palette.navy
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.58), Inches(1.2), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(color)
    line.line.fill.background()


def add_footer(slide, palette, meta: dict, page_num: int):
    """Footer with page number."""
    slide_w_in = meta.get("slide_width", 13.333)
    dark = is_dark(palette)

    # Thin line
    line_width = slide_w_in - 1.2
    line_color = palette.card_border if dark else palette.light_gray
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(7.0), Inches(line_width), Pt(1)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(line_color)
    line.line.fill.background()

    # Left: branding
    lb = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(4), Inches(0.35))
    tf = lb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = meta.get("footer_left", "CROSSINVEST SA | Lugano")
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.color.rgb = rgb(palette.gold)

    # Center: confidential
    cb = slide.shapes.add_textbox(Inches(4.5), Inches(7.05), Inches(4.333), Inches(0.35))
    tf2 = cb.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = meta.get("footer_center", "Strictly Confidential")
    r2.font.size = Pt(8)
    r2.font.italic = True
    r2.font.color.rgb = rgb(palette.source_color)

    # Right: page number
    rb = slide.shapes.add_textbox(Inches(10), Inches(7.05), Inches(2.733), Inches(0.35))
    tf3 = rb.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.RIGHT
    r3 = p3.add_run()
    r3.text = f"Page {page_num}"
    r3.font.size = Pt(8)
    r3.font.color.rgb = rgb(palette.source_color)


def add_slide_title(slide, title: str, subtitle: str = None, tag: str = None, palette=None):
    """Assertion-evidence title block — uses title_color for dark/light awareness."""
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.7), Inches(11.5), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = rgb(palette.title_color)
    run.font.name = "Arial"

    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(1.45), Inches(11.5), Inches(0.5))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(12)
        r2.font.color.rgb = rgb(palette.body_text)

    if tag:
        tb3 = slide.shapes.add_textbox(Inches(0.6), Inches(1.95), Inches(4), Inches(0.3))
        tf3 = tb3.text_frame
        p3 = tf3.paragraphs[0]
        r3 = p3.add_run()
        r3.text = tag
        r3.font.size = Pt(11)
        r3.font.bold = True
        r3.font.color.rgb = rgb(palette.orange)


def add_chart_tag(slide, chart_num: int, left: float = 0.4, top: float = 2.15, palette=None):
    """Chart number badge with glow accent."""
    tag_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(0.9), Inches(0.28)
    )
    tag_shape.fill.solid()
    tag_shape.fill.fore_color.rgb = rgb(palette.orange)
    tag_shape.line.fill.background()
    tf = tag_shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"Chart #{chart_num}"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = rgb(palette.white)


def add_chart_image(slide, img_path: str, left: float = 0.4, top: float = 2.5,
                    width: float = 12.5, height: float = 4.5):
    """Place a chart image on the slide."""
    slide.shapes.add_picture(
        img_path, Inches(left), Inches(top), Inches(width), Inches(height)
    )


def make_content_slide(prs, palette, meta: dict, page_num: int, bg_image=None):
    """Create slide with header, section line, footer. Editorial branch for SYZ-style."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    add_slide_bg(slide, prs, palette, bg_image)

    if is_editorial(palette) or meta.get("mode") == "editorial":
        add_editorial_header(slide, palette, meta)
        add_editorial_footer(slide, palette, meta, page_num)
    else:
        add_header_bar(slide, palette, meta)
        add_section_line(slide, palette)
        add_footer(slide, palette, meta, page_num)
    return slide


# ============================================================================
# Slide Builders
# ============================================================================

def build_cover(prs, content: dict, palette, bg_image=None):
    """
    Full-bleed cover slide with AI-generated background.
    Falls back to solid navy if no bg_image.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Background: AI image or solid navy
    if bg_image and os.path.exists(bg_image):
        slide.shapes.add_picture(bg_image, Emu(0), Emu(0), slide_w, slide_h)
        # Add a dark overlay for text readability
        overlay = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), slide_w, slide_h
        )
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = rgb("#000000")
        overlay.line.fill.background()
        _set_shape_alpha(overlay, 40)  # 40% opaque black overlay
    else:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), slide_w, slide_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = rgb(palette.navy)
        bg.line.fill.background()

    # Gold accent bars
    for y_pos in [0, 7.38]:
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(y_pos), slide_w, Inches(0.12)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(palette.gold)
        bar.line.fill.background()

    elements = [
        (1.0, content.get("company", ""), 38, palette.gold, True),
        (1.7, content.get("tagline", ""), 14, palette.gold, True),
        (2.7, content.get("headline", ""), 44, palette.white, False),
        (3.6, content.get("subheadline", ""), 24, palette.gold, False),
        (4.2, content.get("date_range", ""), 16, "#AABBCC", False),
        (4.85, content.get("theme", ""), 18, palette.white, True),
        (6.4, content.get("footer_text", ""), 11, "#8899AA", False),
    ]

    for y_off, text, size, color, bold in elements:
        if not text:
            continue
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(y_off), Inches(11.7), Inches(0.7))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = "Arial"
        r.font.color.rgb = rgb(color) if isinstance(color, str) and color.startswith("#") else rgb(color)

    # Separator line under tagline
    add_glow_line(slide, 0.8, 2.35, 2.5, palette)


def build_executive_summary(prs, content: dict, palette, meta: dict, page_num: int, bg_image=None):
    """
    SCR-format executive summary with glassmorphism cards.
    """
    slide = make_content_slide(prs, palette, meta, page_num, bg_image)
    add_slide_title(slide, content.get("title", ""), tag=content.get("tag"), palette=palette)

    dark = is_dark(palette)
    y_pos = 2.2

    for section in content.get("sections", []):
        height = section.get("height", 1.15)

        # Add glassmorphism card behind each section
        if dark:
            add_glass_card(slide, 0.4, y_pos - 0.1, 12.4, height, palette, alpha_pct=70)

        # Colored left accent bar for SCR sections
        label = section.get("label", "")
        if label in ("S", "C", "R"):
            accent_colors = {"S": palette.chart_primary, "C": palette.red, "R": palette.green}
            accent = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.45), Inches(y_pos - 0.05),
                Inches(0.08), Inches(height - 0.1)
            )
            accent.fill.solid()
            accent.fill.fore_color.rgb = rgb(accent_colors.get(label, palette.gold))
            accent.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(11.8), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.space_after = Pt(4)
        p.line_spacing = Pt(16)
        run = p.add_run()
        run.text = section.get("text", "")
        run.font.size = Pt(11)
        run.font.color.rgb = rgb(palette.body_text)
        run.font.name = "Arial"
        y_pos += height


def build_chart_slide(prs, content: dict, chart_path: str, chart_num: int,
                      palette, meta: dict, page_num: int, bg_image=None):
    """Chart slide with assertion title + chart image on AI background."""
    slide = make_content_slide(prs, palette, meta, page_num, bg_image)
    add_slide_title(slide, content.get("title", ""), tag=content.get("tag"), palette=palette)
    add_chart_tag(slide, chart_num, palette=palette)

    dims = content.get("chart_dims", {})
    add_chart_image(
        slide, chart_path,
        left=dims.get("left", 0.4),
        top=dims.get("top", 2.5),
        width=dims.get("width", 12.5),
        height=dims.get("height", 4.5),
    )

    # Optional: interactive link button (e.g. to Bloomberg chart, CNBC video)
    link_url = content.get("link_url", "")
    if link_url:
        link_label = content.get("link_label", "📊 Interactive Chart")
        add_link_button(slide, link_url, label=link_label,
                        left=10.5, top=2.15, width=2.2, height=0.32,
                        palette=palette)


def build_text_slide(prs, content: dict, palette, meta: dict, page_num: int, bg_image=None):
    """
    Two-column text slide with glassmorphism cards.
    """
    slide = make_content_slide(prs, palette, meta, page_num, bg_image)
    add_slide_title(slide, content.get("title", ""), tag=content.get("tag"), palette=palette)

    dark = is_dark(palette)

    # Glassmorphism cards for each column
    if dark:
        add_glass_card(slide, 0.4, 2.05, 5.8, 4.6, palette, alpha_pct=65)
        add_glass_card(slide, 6.8, 2.05, 5.8, 4.6, palette, alpha_pct=65)

    for side, x_start in [("left", 0.6), ("right", 7.0)]:
        col_title = content.get(f"{side}_title", "")
        col_color = content.get(f"{side}_color", palette.body_text)
        items = content.get(f"{side}_items", [])
        icon = content.get(f"{side}_icon", "+" if side == "left" else "!")

        if col_title:
            title_box = slide.shapes.add_textbox(
                Inches(x_start), Inches(2.2), Inches(5.5), Inches(0.4)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = col_title
            r.font.size = Pt(14)
            r.font.bold = True
            r.font.color.rgb = rgb(col_color)

        y_item = 2.7
        for item in items:
            tb = slide.shapes.add_textbox(
                Inches(x_start), Inches(y_item), Inches(5.3), Inches(0.45)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.line_spacing = Pt(14)
            run = p.add_run()
            run.text = f"{icon}  {item}"
            run.font.size = Pt(10.5)
            run.font.color.rgb = rgb(palette.body_text)
            y_item += 0.55

    # Vertical separator with glow
    sep_color = palette.card_border if dark else palette.light_gray
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(2.3), Pt(1.5), Inches(4.2)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = rgb(sep_color)
    sep.line.fill.background()


def build_sources_slide(prs, content: dict, palette, meta: dict, page_num: int, bg_image=None):
    """Data sources and disclaimer slide."""
    slide = make_content_slide(prs, palette, meta, page_num, bg_image)
    add_slide_title(slide, content.get("title", "Data Sources & Methodology"), palette=palette)

    dark = is_dark(palette)

    # Glass cards for source columns
    if dark:
        add_glass_card(slide, 0.4, 1.8, 5.8, 4.4, palette, alpha_pct=60)
        add_glass_card(slide, 6.8, 1.8, 5.8, 4.4, palette, alpha_pct=60)

    # Left column
    tb_l = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(5.8), Inches(4.2))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    p_l = tf_l.paragraphs[0]
    p_l.line_spacing = Pt(13)
    r_l = p_l.add_run()
    r_l.text = content.get("left_sources", "")
    r_l.font.size = Pt(8)
    r_l.font.color.rgb = rgb(palette.body_text)

    # Right column
    tb_r = slide.shapes.add_textbox(Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.2))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.line_spacing = Pt(13)
    r_r = p_r.add_run()
    r_r.text = content.get("right_sources", "")
    r_r.font.size = Pt(8)
    r_r.font.color.rgb = rgb(palette.body_text)

    # Separator
    sep_color = palette.card_border if dark else palette.light_gray
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(2.1), Pt(1), Inches(3.8)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = rgb(sep_color)
    sep.line.fill.background()

    # Disclaimer
    disclaimer = content.get("disclaimer", "")
    if disclaimer:
        tb_d = slide.shapes.add_textbox(Inches(0.6), Inches(6.2), Inches(12.0), Inches(0.7))
        tf_d = tb_d.text_frame
        tf_d.word_wrap = True
        p_d = tf_d.paragraphs[0]
        p_d.line_spacing = Pt(11)
        r_d = p_d.add_run()
        r_d.text = disclaimer
        r_d.font.size = Pt(7)
        r_d.font.italic = True
        r_d.font.color.rgb = rgb(palette.source_color)


def build_back_cover(prs, content: dict, palette, bg_image=None):
    """Navy back cover with contact details and optional AI background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Background
    if bg_image and os.path.exists(bg_image):
        slide.shapes.add_picture(bg_image, Emu(0), Emu(0), slide_w, slide_h)
        overlay = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), slide_w, slide_h
        )
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = rgb("#000000")
        overlay.line.fill.background()
        _set_shape_alpha(overlay, 50)
    else:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), slide_w, slide_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = rgb(palette.navy)
        bg.line.fill.background()

    # Gold accent lines
    for y_pos in [2.0, 5.5]:
        add_glow_line(slide, 4.5, y_pos, 4.333, palette)

    # Company name
    tb1 = slide.shapes.add_textbox(Inches(0), Inches(2.3), prs.slide_width, Inches(0.8))
    tf1 = tb1.text_frame
    p1 = tf1.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = content.get("company", "CROSSINVEST SA")
    r1.font.size = Pt(36)
    r1.font.bold = True
    r1.font.color.rgb = rgb(palette.gold)
    r1.font.name = "Arial"

    # Tagline
    tb2 = slide.shapes.add_textbox(Inches(0), Inches(3.1), prs.slide_width, Inches(0.5))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = content.get("tagline", "")
    r2.font.size = Pt(16)
    r2.font.color.rgb = rgb(palette.white)
    r2.font.name = "Arial"

    # Contact lines
    contact_y = 3.8
    for line_text in content.get("contact_lines", []):
        tb = slide.shapes.add_textbox(Inches(0), Inches(contact_y), prs.slide_width, Inches(0.35))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = line_text
        r.font.size = Pt(12)
        r.font.color.rgb = rgb("#AABBCC")
        r.font.name = "Arial"
        contact_y += 0.35

    # Closing message
    closing = content.get("closing", "")
    if closing:
        tb3 = slide.shapes.add_textbox(Inches(0), Inches(5.8), prs.slide_width, Inches(0.5))
        tf3 = tb3.text_frame
        p3 = tf3.paragraphs[0]
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = closing
        r3.font.size = Pt(14)
        r3.font.italic = True
        r3.font.color.rgb = rgb(palette.gold)

    # Regulatory & copyright
    for y_off, text_key, size in [(6.4, "regulatory", 9), (6.8, "copyright", 8)]:
        text = content.get(text_key, "")
        if text:
            tb = slide.shapes.add_textbox(Inches(0), Inches(y_off), prs.slide_width, Inches(0.35))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.color.rgb = rgb("#8899AA")


def build_section_divider(prs, content: dict, palette, bg_image=None, meta=None):
    """
    Full-bleed section divider slide.
    Editorial mode: white bg + colored accent bar on left + title in section color.
    Standard mode: AI background or solid navy with gold accents.

    content:
      section_num: "01"
      title: "GLOBAL MACRO"
      subtitle: "Economic Indicators, Central Banks & Growth Outlook"
      section: "macro"  (for editorial badge color lookup)
    """
    meta = meta or {}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    if is_editorial(palette) or meta.get("mode") == "editorial":
        # Editorial mode: white background + colored accent bar
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), slide_w, slide_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = rgb("#FFFFFF")
        bg.line.fill.background()

        section = content.get("section", "")
        accent_color = SECTION_BADGE_COLORS.get(section, palette.orange)

        # Thick colored accent bar on left
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            Inches(0.15), slide_h
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(accent_color)
        bar.line.fill.background()

        # Section badge
        if section:
            add_section_badge(slide, section, 0.6, 2.5, palette)

        # Section title in section color
        title = content.get("title", "")
        if title:
            tb_t = slide.shapes.add_textbox(Inches(0.6), Inches(3.0), Inches(11), Inches(1.2))
            tf = tb_t.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = title
            r.font.size = Pt(40)
            r.font.bold = True
            r.font.color.rgb = rgb(accent_color)
            r.font.name = "Arial"

        # Subtitle
        subtitle = content.get("subtitle", content.get("tag", ""))
        if subtitle:
            tb_s = slide.shapes.add_textbox(Inches(0.6), Inches(4.3), Inches(10), Inches(0.8))
            tf = tb_s.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = subtitle
            r.font.size = Pt(14)
            r.font.color.rgb = rgb(palette.gray)
            r.font.name = "Arial"

        # Thin accent line below subtitle
        add_glow_line(slide, 0.6, 5.2, 3.0, palette, color=accent_color)
    else:
        # Standard mode: AI background or solid navy
        if bg_image and os.path.exists(bg_image):
            slide.shapes.add_picture(bg_image, Emu(0), Emu(0), slide_w, slide_h)
            overlay = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), slide_w, slide_h
            )
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = rgb("#000000")
            overlay.line.fill.background()
            _set_shape_alpha(overlay, 45)
        else:
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), slide_w, slide_h)
            bg.fill.solid()
            bg.fill.fore_color.rgb = rgb(palette.navy)
            bg.line.fill.background()

        # Gold accent line
        add_glow_line(slide, 1.0, 2.8, 3.0, palette)

        # Section number (large, gold)
        num = content.get("section_num", "")
        if num:
            tb_num = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(3), Inches(1.0))
            tf = tb_num.text_frame
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = num
            r.font.size = Pt(72)
            r.font.bold = True
            r.font.color.rgb = rgb(palette.gold)
            r.font.name = "Arial"

        # Section title (large, white)
        title = content.get("title", "")
        if title:
            tb_t = slide.shapes.add_textbox(Inches(1.0), Inches(3.1), Inches(11), Inches(1.2))
            tf = tb_t.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = title
            r.font.size = Pt(44)
            r.font.bold = True
            r.font.color.rgb = rgb(palette.white)
            r.font.name = "Arial"

        # Subtitle
        subtitle = content.get("subtitle", "")
        if subtitle:
            tb_s = slide.shapes.add_textbox(Inches(1.0), Inches(4.4), Inches(10), Inches(0.8))
            tf = tb_s.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = subtitle
            r.font.size = Pt(16)
            r.font.color.rgb = rgb("#AABBCC")
            r.font.name = "Arial"

        # Bottom gold line
        add_glow_line(slide, 1.0, 5.5, 11.3, palette)

        # Optional video link for section intro
        video_url = content.get("video_url", "")
        if video_url:
            video_label = content.get("video_label", "▶ Section Overview")
            add_video_button(slide, video_url, label=video_label,
                             left=1.0, top=5.8, width=2.5, height=0.4,
                             palette=palette)

        # Company branding bottom-right
        company = meta.get("company", "Rossignoli & Partners")
        tb_brand = slide.shapes.add_textbox(Inches(9), Inches(6.6), Inches(4), Inches(0.35))
        tf = tb_brand.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = f"{company}  |  Weekly Strategy"
        r.font.size = Pt(9)
        r.font.color.rgb = rgb(palette.gold)


def build_kpi_dashboard(prs, content: dict, palette, meta: dict, page_num: int, bg_image=None):
    """
    KPI Dashboard with glass cards showing key metrics.

    content:
      title: "Market Dashboard"
      tag: "WEEKLY SNAPSHOT"
      kpis: [
        {"label": "DOW JONES", "value": "50,115", "change": "+2.5%", "positive": true},
        {"label": "S&P 500", "value": "6,932", "change": "+1.2%", "positive": true},
        ...
      ]
    """
    slide = make_content_slide(prs, palette, meta, page_num, bg_image)
    add_slide_title(slide, content.get("title", "Market Dashboard"),
                    tag=content.get("tag"), palette=palette)

    dark = is_dark(palette)
    kpis = content.get("kpis", [])

    # Calculate grid layout: max 4 per row
    cols = min(len(kpis), 4)
    rows_count = (len(kpis) + cols - 1) // cols

    card_w = 2.8
    card_h = 1.8
    gap = 0.25
    total_w = cols * card_w + (cols - 1) * gap
    start_x = (13.333 - total_w) / 2
    start_y = 2.3

    for idx, kpi in enumerate(kpis):
        row = idx // cols
        col = idx % cols
        x = start_x + col * (card_w + gap)
        y = start_y + row * (card_h + gap)

        # Glass card
        add_glass_card(slide, x, y, card_w, card_h, palette,
                       alpha_pct=70 if dark else 100)

        # Label (small, gold)
        tb_label = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 0.15), Inches(card_w - 0.4), Inches(0.3)
        )
        tf = tb_label.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = kpi.get("label", "")
        r.font.size = Pt(9)
        r.font.bold = True
        r.font.color.rgb = rgb(palette.gold)

        # Hero value (large, white)
        tb_val = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 0.45), Inches(card_w - 0.4), Inches(0.7)
        )
        tf = tb_val.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = kpi.get("value", "")
        r.font.size = Pt(28)
        r.font.bold = True
        r.font.color.rgb = rgb(palette.title_color)
        r.font.name = "Arial"

        # Change (color-coded)
        change = kpi.get("change", "")
        positive = kpi.get("positive", True)
        tb_chg = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 1.2), Inches(card_w - 0.4), Inches(0.35)
        )
        tf = tb_chg.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = change
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = rgb(palette.green if positive else palette.red)


def build_news_slide(prs, content: dict, palette, meta: dict, page_num: int, bg_image=None):
    """
    News slide with headline, source, summary, and market impact.

    content:
      title: "Key News This Week"
      stories: [
        {
          "headline": "Trump Announces 25% Tariffs on EU Steel",
          "source": "Reuters | Feb 5, 2026",
          "summary": "President Trump...",
          "impact": "Markets dropped 1.2% on the announcement...",
          "sentiment": "negative"  # positive, negative, neutral
        },
        ...
      ]
    """
    slide = make_content_slide(prs, palette, meta, page_num, bg_image)
    add_slide_title(slide, content.get("title", "Key News"),
                    tag=content.get("tag"), palette=palette)

    dark = is_dark(palette)
    stories = content.get("stories", [])
    y_pos = 2.2

    for story in stories[:3]:  # max 3 stories per slide
        height = 1.5

        # Glass card for each story
        if dark:
            add_glass_card(slide, 0.4, y_pos - 0.05, 12.4, height, palette, alpha_pct=65)

        # Sentiment indicator bar (left edge)
        sentiment = story.get("sentiment", "neutral")
        sent_colors = {"positive": palette.green, "negative": palette.red, "neutral": palette.orange}
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.45), Inches(y_pos),
            Inches(0.06), Inches(height - 0.1)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = rgb(sent_colors.get(sentiment, palette.orange))
        accent.line.fill.background()

        # Headline (bold, large — clickable if url provided)
        tb_h = slide.shapes.add_textbox(
            Inches(0.7), Inches(y_pos + 0.05), Inches(11.8), Inches(0.35)
        )
        tf = tb_h.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]

        headline_url = story.get("url", "")
        if headline_url:
            add_hyperlink_run(p, story.get("headline", ""), headline_url,
                              font_size=13, color=palette.title_color, bold=True)
        else:
            r = p.add_run()
            r.text = story.get("headline", "")
            r.font.size = Pt(13)
            r.font.bold = True
            r.font.color.rgb = rgb(palette.title_color)

        # Video link button (if video_url provided)
        video_url = story.get("video_url", "")
        if video_url:
            add_video_button(slide, video_url, label="▶ Watch",
                             left=11.5, top=y_pos + 0.05, width=1.3, height=0.32,
                             palette=palette)

        # Source (small, gray — clickable if source_url provided)
        tb_src = slide.shapes.add_textbox(
            Inches(0.7), Inches(y_pos + 0.4), Inches(6), Inches(0.25)
        )
        tf = tb_src.text_frame
        p = tf.paragraphs[0]
        source_url = story.get("source_url", "")
        if source_url:
            add_hyperlink_run(p, story.get("source", ""), source_url,
                              font_size=8, color=palette.source_color, italic=True)
        else:
            r = p.add_run()
            r.text = story.get("source", "")
            r.font.size = Pt(8)
            r.font.italic = True
            r.font.color.rgb = rgb(palette.source_color)

        # Summary + Impact (two columns)
        summary = story.get("summary", "")
        impact = story.get("impact", "")

        if summary:
            tb_sum = slide.shapes.add_textbox(
                Inches(0.7), Inches(y_pos + 0.7), Inches(6.5), Inches(0.65)
            )
            tf = tb_sum.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.line_spacing = Pt(13)
            r = p.add_run()
            r.text = summary
            r.font.size = Pt(9)
            r.font.color.rgb = rgb(palette.body_text)

        if impact:
            tb_imp = slide.shapes.add_textbox(
                Inches(7.5), Inches(y_pos + 0.7), Inches(5.2), Inches(0.65)
            )
            tf = tb_imp.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.line_spacing = Pt(13)
            r = p.add_run()
            r.text = impact
            r.font.size = Pt(9)
            r.font.bold = True
            r.font.color.rgb = rgb(sent_colors.get(sentiment, palette.body_text))

        y_pos += height + 0.1


def build_image_slide(prs, content: dict, palette, meta: dict, page_num: int, bg_image=None):
    """
    Embed a real image (web screenshot, chart, infographic) on a slide.
    The image is the hero element — full width with optional caption and source.

    content:
      title: "S&P 500 Year-to-Date Performance"
      tag: "REAL-TIME DATA"
      image_path: "/tmp/real_charts/sp500_ytd.png"
      caption: "Source: TradingView | Data as of Feb 7, 2026"
      source_url: "https://www.tradingview.com/..."
      commentary: "Optional analysis text below the image"
    """
    slide = make_content_slide(prs, palette, meta, page_num, bg_image)
    add_slide_title(slide, content.get("title", ""), tag=content.get("tag"), palette=palette)

    dark = is_dark(palette)
    image_path = content.get("image_path", "")

    if image_path and os.path.exists(image_path):
        # Glass card frame around the image
        if dark:
            add_glass_card(slide, 0.3, 2.1, 12.7, 4.2, palette, alpha_pct=50)

        # Embed image — centered, max width
        img_left = content.get("img_left", 0.4)
        img_top = content.get("img_top", 2.2)
        img_width = content.get("img_width", 12.5)
        img_height = content.get("img_height", 3.9)
        slide.shapes.add_picture(
            image_path,
            Inches(img_left), Inches(img_top),
            Inches(img_width), Inches(img_height)
        )
    else:
        # Fallback: show a "no image" message
        tb = slide.shapes.add_textbox(Inches(3), Inches(3.5), Inches(7), Inches(1))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"[Image not found: {image_path}]"
        r.font.size = Pt(14)
        r.font.color.rgb = rgb(palette.source_color)

    # Caption / source attribution
    caption = content.get("caption", "")
    if caption:
        tb_cap = slide.shapes.add_textbox(Inches(0.6), Inches(6.3), Inches(8), Inches(0.3))
        tf = tb_cap.text_frame
        p = tf.paragraphs[0]
        source_url = content.get("source_url", "")
        if source_url:
            add_hyperlink_run(p, caption, source_url,
                              font_size=8, color=palette.source_color, italic=True)
        else:
            r = p.add_run()
            r.text = caption
            r.font.size = Pt(8)
            r.font.italic = True
            r.font.color.rgb = rgb(palette.source_color)

    # Optional commentary text
    commentary = content.get("commentary", "")
    if commentary:
        if dark:
            add_glass_card(slide, 0.3, 6.1, 12.7, 0.7, palette, alpha_pct=60)
        tb_comm = slide.shapes.add_textbox(Inches(0.6), Inches(6.15), Inches(12), Inches(0.6))
        tf = tb_comm.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = commentary
        r.font.size = Pt(10)
        r.font.color.rgb = rgb(palette.body_text)

    # Link button if source_url provided
    link_url = content.get("source_url", "")
    if link_url:
        add_link_button(slide, link_url, label="📊 Live Chart",
                        left=10.5, top=6.3, width=2.2, height=0.32,
                        palette=palette)


def build_dual_chart_slide(prs, content: dict, chart_paths: list, palette,
                           meta: dict, page_num: int, bg_image=None):
    """
    Side-by-side dual chart comparison slide.

    content:
      title: "Gold vs Bitcoin: Diverging Safe Havens"
      tag: "CROSS-ASSET"
      left_label: "Gold: +14.4% YTD"
      right_label: "Bitcoin: -20.5% YTD"
    chart_paths: ["/path/to/left_chart.png", "/path/to/right_chart.png"]
    """
    slide = make_content_slide(prs, palette, meta, page_num, bg_image)
    add_slide_title(slide, content.get("title", ""), tag=content.get("tag"), palette=palette)

    dark = is_dark(palette)

    # Left chart
    if len(chart_paths) > 0 and os.path.exists(chart_paths[0]):
        if dark:
            add_glass_card(slide, 0.3, 2.3, 6.2, 4.0, palette, alpha_pct=50)
        slide.shapes.add_picture(
            chart_paths[0], Inches(0.4), Inches(2.4), Inches(6.0), Inches(3.7)
        )

    # Right chart
    if len(chart_paths) > 1 and os.path.exists(chart_paths[1]):
        if dark:
            add_glass_card(slide, 6.8, 2.3, 6.2, 4.0, palette, alpha_pct=50)
        slide.shapes.add_picture(
            chart_paths[1], Inches(6.9), Inches(2.4), Inches(6.0), Inches(3.7)
        )

    # Labels
    for i, (x_start, key) in enumerate([(0.4, "left_label"), (6.9, "right_label")]):
        label = content.get(key, "")
        if label:
            tb = slide.shapes.add_textbox(
                Inches(x_start), Inches(2.05), Inches(6.0), Inches(0.3)
            )
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = label
            r.font.size = Pt(11)
            r.font.bold = True
            r.font.color.rgb = rgb(palette.gold)

    # Vertical separator
    sep_color = palette.gold if dark else palette.navy
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.55), Inches(2.5), Pt(2), Inches(3.5)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = rgb(sep_color)
    sep.line.fill.background()


def build_callout_slide(prs, content: dict, palette, meta: dict, page_num: int, bg_image=None):
    """
    Big callout slide — one hero number/quote with context.
    JPMorgan "Guide to Markets" style: huge number + explanation.

    content:
      title: "Key Metric"
      tag: "WEEKLY HIGHLIGHT"
      hero_number: "$4,965"
      hero_label: "Gold Price ($/oz)"
      hero_color: "#F0B90B"  (optional, default gold)
      context_left: "JPM Target: $6,300"
      context_right: "+14.4% YTD"
      commentary: "Gold's best start to a year since 1980..."
    """
    slide = make_content_slide(prs, palette, meta, page_num, bg_image)
    add_slide_title(slide, content.get("title", ""), tag=content.get("tag"), palette=palette)

    dark = is_dark(palette)
    hero_color = content.get("hero_color", palette.gold)

    # Glass card for hero section
    if dark:
        add_glass_card(slide, 1.5, 2.2, 10.3, 3.8, palette, alpha_pct=65)

    # Hero number (massive)
    hero = content.get("hero_number", "")
    if hero:
        tb_hero = slide.shapes.add_textbox(
            Inches(0), Inches(2.4), Inches(13.333), Inches(1.5)
        )
        tf = tb_hero.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = hero
        r.font.size = Pt(72)
        r.font.bold = True
        r.font.color.rgb = rgb(hero_color)
        r.font.name = "Arial"

    # Hero label
    hero_label = content.get("hero_label", "")
    if hero_label:
        tb_label = slide.shapes.add_textbox(
            Inches(0), Inches(3.8), Inches(13.333), Inches(0.5)
        )
        tf = tb_label.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = hero_label
        r.font.size = Pt(18)
        r.font.color.rgb = rgb(palette.body_text)

    # Context: left and right metrics
    ctx_left = content.get("context_left", "")
    ctx_right = content.get("context_right", "")

    if ctx_left:
        tb_cl = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(4.5), Inches(0.5))
        tf = tb_cl.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = ctx_left
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = rgb(palette.chart_primary)

    if ctx_right:
        tb_cr = slide.shapes.add_textbox(Inches(6.8), Inches(4.5), Inches(4.5), Inches(0.5))
        tf = tb_cr.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = ctx_right
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = rgb(palette.green)

    # Glow accent line
    add_glow_line(slide, 3, 4.35, 7.3, palette, color=hero_color)

    # Commentary text
    commentary = content.get("commentary", "")
    if commentary:
        if dark:
            add_glass_card(slide, 1.5, 5.2, 10.3, 1.2, palette, alpha_pct=55)
        tb_comm = slide.shapes.add_textbox(Inches(2), Inches(5.3), Inches(9.3), Inches(1.0))
        tf = tb_comm.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = Pt(16)
        r = p.add_run()
        r.text = commentary
        r.font.size = Pt(12)
        r.font.color.rgb = rgb(palette.body_text)


# ============================================================================
# Editorial Slide Builders (SYZ-style)
# ============================================================================

def build_editorial(prs, content: dict, palette, meta: dict, page_num: int,
                    chart_path: str = None, bg_image=None):
    """
    SYZ-style editorial slide: badge + hashtags + commentary + chart/image + source.

    content:
      section: "equities"
      hashtags: "#us #equities #sp500 #technicals"
      commentary: "The S&P 500 extended its rally..."
      image_path: "/path/to/screenshot.png" (alternative to chart_path)
      source: "Source: Bloomberg"
      source_url: "https://..."
      title: "Optional override title"
      chart_dims: {left, top, width, height}
    """
    slide = make_content_slide(prs, palette, meta, page_num, bg_image)
    el = EDITORIAL_LAYOUT
    section = content.get("section", "")
    section_color = SECTION_BADGE_COLORS.get(section, palette.orange)

    # Section badge
    if section:
        add_section_badge(slide, section, el.margin_left, el.badge_top, palette)

    # Hashtag title in section color
    hashtags = content.get("hashtags", "")
    if hashtags:
        tb_hash = slide.shapes.add_textbox(
            Inches(el.margin_left), Inches(el.hashtag_top),
            Inches(el.content_width), Inches(0.3)
        )
        tf = tb_hash.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = hashtags
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = rgb(section_color)
        r.font.name = "Arial"

    # Title (if provided, goes above commentary)
    title = content.get("title", "")
    if title:
        tb_title = slide.shapes.add_textbox(
            Inches(el.margin_left), Inches(el.commentary_top - 0.35),
            Inches(el.content_width), Inches(0.35)
        )
        tf = tb_title.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = rgb(palette.body_text)
        r.font.name = "Arial"

    # Editorial commentary
    commentary = content.get("commentary", "")
    if commentary:
        comm_top = el.commentary_top if not title else el.commentary_top + 0.1
        tb_comm = slide.shapes.add_textbox(
            Inches(el.margin_left), Inches(comm_top),
            Inches(el.content_width), Inches(el.commentary_height)
        )
        tf = tb_comm.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.line_spacing = Pt(16)
        r = p.add_run()
        r.text = commentary
        r.font.size = Pt(11)
        r.font.color.rgb = rgb(palette.body_text)
        r.font.name = "Arial"

    # Chart image or external image
    dims = content.get("chart_dims", {})
    img_left = dims.get("left", el.margin_left)
    img_top = dims.get("top", el.chart_top)
    img_width = dims.get("width", el.content_width)
    img_height = dims.get("height", el.chart_height)

    img_path = chart_path or content.get("image_path", "")
    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(
            img_path,
            Inches(img_left), Inches(img_top),
            Inches(img_width), Inches(img_height)
        )
    elif content.get("image_path"):
        # Fallback placeholder
        tb = slide.shapes.add_textbox(
            Inches(img_left + 2), Inches(img_top + 1),
            Inches(8), Inches(1)
        )
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"[Image: {content['image_path']}]"
        r.font.size = Pt(12)
        r.font.color.rgb = rgb(palette.source_color)

    # Source attribution (right-aligned, gray italic)
    source = content.get("source", "")
    if source:
        tb_src = slide.shapes.add_textbox(
            Inches(el.margin_left), Inches(el.source_top),
            Inches(el.content_width), Inches(0.3)
        )
        tf = tb_src.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        source_url = content.get("source_url", "")
        if source_url:
            add_hyperlink_run(p, source, source_url,
                              font_size=8, color=palette.source_color, italic=True)
        else:
            r = p.add_run()
            r.text = source
            r.font.size = Pt(8)
            r.font.italic = True
            r.font.color.rgb = rgb(palette.source_color)


def build_quote_slide(prs, content: dict, palette, meta: dict, page_num: int, bg_image=None):
    """
    CEO/analyst quote slide with attribution.

    content:
      quote: "The markets are pricing in..."
      attribution: "Jamie Dimon, CEO JPMorgan"
      source: "Bloomberg Interview, Feb 2026"
      section: "macro"
      highlight: true  (yellow highlight box)
      commentary: "Optional brief context below"
    """
    slide = make_content_slide(prs, palette, meta, page_num, bg_image)
    section = content.get("section", "")
    section_color = SECTION_BADGE_COLORS.get(section, palette.orange)

    # Section badge
    if section:
        add_section_badge(slide, section, 0.6, 0.7, palette)

    # Optional yellow highlight box
    if content.get("highlight", False):
        highlight = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.5), Inches(1.8), Inches(10.3), Inches(3.0)
        )
        highlight.fill.solid()
        highlight.fill.fore_color.rgb = rgb("#FFF8E1")
        highlight.line.color.rgb = rgb("#FFE082")
        highlight.line.width = Pt(1)

    # Large curly quote
    quote_text = content.get("quote", "")
    if quote_text:
        tb_q = slide.shapes.add_textbox(
            Inches(2.0), Inches(2.0), Inches(9.3), Inches(2.2)
        )
        tf = tb_q.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = Pt(24)
        r = p.add_run()
        r.text = f"\u201C{quote_text}\u201D"
        r.font.size = Pt(20)
        r.font.italic = True
        r.font.color.rgb = rgb(palette.body_text)
        r.font.name = "Georgia"

    # Attribution with em-dash
    attribution = content.get("attribution", "")
    if attribution:
        tb_attr = slide.shapes.add_textbox(
            Inches(2.0), Inches(4.4), Inches(9.3), Inches(0.4)
        )
        tf = tb_attr.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = f"\u2014 {attribution}"
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = rgb(section_color)
        r.font.name = "Arial"

    # Source
    source = content.get("source", "")
    if source:
        tb_src = slide.shapes.add_textbox(
            Inches(2.0), Inches(4.85), Inches(9.3), Inches(0.3)
        )
        tf = tb_src.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = source
        r.font.size = Pt(9)
        r.font.italic = True
        r.font.color.rgb = rgb(palette.source_color)

    # Commentary
    commentary = content.get("commentary", "")
    if commentary:
        tb_comm = slide.shapes.add_textbox(
            Inches(1.5), Inches(5.5), Inches(10.3), Inches(1.0)
        )
        tf = tb_comm.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.line_spacing = Pt(16)
        r = p.add_run()
        r.text = commentary
        r.font.size = Pt(11)
        r.font.color.rgb = rgb(palette.body_text)


def build_chart_grid(prs, content: dict, palette, meta: dict, page_num: int,
                     chart_paths: list = None, bg_image=None):
    """
    2x2 or 2x3 grid of charts/images for comparison.

    content:
      title: "Precious Metals Dashboard"
      section: "commodities"
      hashtags: "#gold #silver #platinum #palladium"
      grid: [
        {"label": "Gold YTD", "image_path": "/path/to/gold.png"},
        {"label": "Silver YTD", "image_path": "/path/to/silver.png"},
        ...
      ]
      cols: 2  (default 2)
      source: "Source: Bloomberg"
    """
    slide = make_content_slide(prs, palette, meta, page_num, bg_image)
    section = content.get("section", "")
    section_color = SECTION_BADGE_COLORS.get(section, palette.orange)

    # Section badge
    if section:
        add_section_badge(slide, section, 0.6, 0.7, palette)

    # Title
    title = content.get("title", "")
    if title:
        tb_t = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(11.5), Inches(0.4))
        tf = tb_t.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = rgb(palette.body_text)
        r.font.name = "Arial"

    # Hashtags
    hashtags = content.get("hashtags", "")
    if hashtags:
        tb_h = slide.shapes.add_textbox(Inches(0.6), Inches(1.45), Inches(11.5), Inches(0.3))
        tf = tb_h.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = hashtags
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = rgb(section_color)

    # Grid layout
    grid = content.get("grid", [])
    paths = chart_paths or []
    cols = content.get("cols", 2)
    n_items = max(len(grid), len(paths))
    if n_items == 0:
        return

    rows = (n_items + cols - 1) // cols
    margin_l = 0.5
    margin_r = 0.5
    gap = 0.3
    grid_top = 1.9
    grid_bottom = 6.4
    total_w = 13.333 - margin_l - margin_r
    total_h = grid_bottom - grid_top
    cell_w = (total_w - (cols - 1) * gap) / cols
    cell_h = (total_h - (rows - 1) * gap) / rows

    for idx in range(n_items):
        row = idx // cols
        col = idx % cols
        x = margin_l + col * (cell_w + gap)
        y = grid_top + row * (cell_h + gap)

        # Get image path: from chart_paths list or grid item
        img_path = ""
        label = ""
        if idx < len(grid):
            item = grid[idx]
            img_path = item.get("image_path", "")
            label = item.get("label", "")
        if idx < len(paths) and paths[idx]:
            img_path = paths[idx]

        if img_path and os.path.exists(img_path):
            # Leave room for label below image
            label_h = 0.3 if label else 0
            slide.shapes.add_picture(
                img_path,
                Inches(x), Inches(y),
                Inches(cell_w), Inches(cell_h - label_h)
            )

        if label:
            tb_lbl = slide.shapes.add_textbox(
                Inches(x), Inches(y + cell_h - 0.3),
                Inches(cell_w), Inches(0.3)
            )
            tf = tb_lbl.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = label
            r.font.size = Pt(9)
            r.font.bold = True
            r.font.color.rgb = rgb(palette.body_text)

    # Source
    source = content.get("source", "")
    if source:
        tb_src = slide.shapes.add_textbox(
            Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.3)
        )
        tf = tb_src.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = source
        r.font.size = Pt(8)
        r.font.italic = True
        r.font.color.rgb = rgb(palette.source_color)

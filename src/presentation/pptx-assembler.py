#!/usr/bin/env python3
"""
Genesis Presentation Engine — PPTX Assembler

Takes a list of PNG slide screenshots and assembles them into a full-bleed PPTX.
Each screenshot becomes one slide, stretched to fill the entire slide area.

Usage:
  echo '{"screenshots": ["/tmp/s/slide_000.png", ...], "output_path": "/tmp/out.pptx"}' | python3 pptx-assembler.py

Output (stdout):
  {"success": true, "path": "/tmp/out.pptx", "slides": 15}
"""

import json
import sys
import os

from pptx import Presentation
from pptx.util import Inches, Emu


def assemble(screenshots: list, output_path: str, width: float = 13.333, height: float = 7.5) -> dict:
    """Assemble PNG screenshots into a full-bleed PPTX."""
    prs = Presentation()
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)

    slide_count = 0
    for png_path in screenshots:
        if not os.path.exists(png_path):
            print(f"WARNING: Screenshot not found: {png_path}", file=sys.stderr)
            continue

        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.add_picture(
            png_path,
            Emu(0), Emu(0),
            prs.slide_width, prs.slide_height,
        )
        slide_count += 1

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)

    return {
        "success": True,
        "path": output_path,
        "slides": slide_count,
    }


def main():
    try:
        raw = sys.stdin.read()
        if not raw.strip():
            print(json.dumps({"success": False, "error": "Empty input", "slides": 0}))
            sys.exit(1)

        spec = json.loads(raw)
        screenshots = spec.get("screenshots", [])
        output_path = spec.get("output_path", "/tmp/assembled.pptx")
        width = spec.get("width", 13.333)
        height = spec.get("height", 7.5)

        result = assemble(screenshots, output_path, width, height)
        print(json.dumps(result))

    except Exception as e:
        import traceback
        traceback.print_exc(file=sys.stderr)
        print(json.dumps({"success": False, "error": str(e), "slides": 0}))
        sys.exit(1)


if __name__ == "__main__":
    main()
